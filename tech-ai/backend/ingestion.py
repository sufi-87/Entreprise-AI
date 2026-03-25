import os
import io
import csv
import json
import shutil
import zipfile
import fitz
import docx
import pandas as pd
from PIL import Image
from backend.config import DOCS_DIR
from backend.database import get_db_status, upsert_document, list_documents as list_documents_db, delete_document_record
from typing import List, Dict

try:
    import pytesseract
except Exception:
    pytesseract = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


SUPPORTED_EXTENSIONS = {
    'pdf', 'doc', 'docx', 'xls', 'xlsx', 'csv', 'txt', 'ppt', 'pptx',
    'png', 'jpg', 'jpeg', 'webp', 'bmp', 'tif', 'tiff', 'md', 'json', 'xml', 'html', 'htm', 'rtf'
}


def save_uploaded_file(file_obj, plant: str, filename: str) -> str:
    plant_dir = os.path.join(DOCS_DIR, plant)
    os.makedirs(plant_dir, exist_ok=True)
    file_path = os.path.join(plant_dir, filename)
    with open(file_path, 'wb') as buffer:
        shutil.copyfileobj(file_obj.file, buffer)

    stat = os.stat(file_path)
    if get_db_status()['enabled']:
        try:
            upsert_document(plant, filename, filename.split('.')[-1], file_path, stat.st_size, int(stat.st_mtime * 1000))
        except Exception as e:
            print(f'DB document upsert failed: {e}')
    return file_path


def delete_document(plant: str, filename: str) -> bool:
    file_path = os.path.join(DOCS_DIR, plant, filename)
    deleted = False
    if os.path.exists(file_path):
        os.remove(file_path)
        deleted = True
    if get_db_status()['enabled']:
        try:
            delete_document_record(plant, filename)
        except Exception as e:
            print(f'DB document delete failed: {e}')
    return deleted


def get_documents(plant: str = None) -> List[Dict]:
    if get_db_status()['enabled']:
        try:
            docs = list_documents_db(plant)
            if docs:
                return docs
        except Exception as e:
            print(f'DB document list failed: {e}')

    docs = []
    plants_to_check = [plant] if plant else os.listdir(DOCS_DIR)
    for p in plants_to_check:
        plant_path = os.path.join(DOCS_DIR, p)
        if not os.path.isdir(plant_path):
            continue
        for f in os.listdir(plant_path):
            file_path = os.path.join(plant_path, f)
            if os.path.isfile(file_path):
                stat = os.stat(file_path)
                docs.append({
                    'plant': p,
                    'filename': f,
                    'type': f.split('.')[-1].upper(),
                    'size_bytes': stat.st_size,
                    'upload_date': int(stat.st_mtime * 1000),
                    'file_path': file_path,
                })
    return sorted(docs, key=lambda x: x['upload_date'], reverse=True)


def ocr_image(image: Image.Image) -> str:
    if pytesseract is None:
        return ''
    try:
        return pytesseract.image_to_string(image)
    except Exception as e:
        print(f'OCR image failed: {e}')
        return ''


def parse_pdf(file_path: str) -> List[Dict]:
    segments = []
    try:
        doc = fitz.open(file_path)
        for page_idx, page in enumerate(doc):
            page_text = page.get_text('text') or ''
            if not page_text.strip():
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                image = Image.open(io.BytesIO(pix.tobytes('png')))
                page_text = ocr_image(image)
            if page_text.strip():
                segments.append({'text': page_text, 'page_num': page_idx + 1, 'source_type': 'pdf_page'})

            # Extract page images for OCR/drawing labels
            try:
                for img_idx, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base = doc.extract_image(xref)
                    image = Image.open(io.BytesIO(base['image']))
                    image_text = ocr_image(image)
                    if image_text.strip():
                        segments.append({
                            'text': f'Image OCR / drawing annotations from page {page_idx+1}:\n{image_text}',
                            'page_num': page_idx + 1,
                            'source_type': 'pdf_image',
                            'extra': {'image_index': img_idx}
                        })
            except Exception as e:
                print(f'PDF image OCR failed: {e}')
    except Exception as e:
        print(f'Error parsing PDF {file_path}: {e}')
    return segments


def parse_docx(file_path: str) -> List[Dict]:
    segments = []
    try:
        doc = docx.Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if vals:
                    text_parts.append(' | '.join(vals))
        if text_parts:
            segments.append({'text': '\n'.join(text_parts), 'source_type': 'docx'})
    except Exception as e:
        print(f'Error parsing DOCX {file_path}: {e}')
    return segments


def parse_excel_csv(file_path: str) -> List[Dict]:
    segments = []
    try:
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
            segments.append({'text': df.to_string(index=False), 'source_type': 'csv'})
        else:
            xls = pd.ExcelFile(file_path)
            for sheet in xls.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                segments.append({'text': f'Sheet: {sheet}\n{df.to_string(index=False)}', 'source_type': 'excel_sheet', 'extra': {'sheet': sheet}})
    except Exception as e:
        print(f'Error parsing Excel/CSV {file_path}: {e}')
    return segments


def parse_txt(file_path: str) -> List[Dict]:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return [{'text': f.read(), 'source_type': 'text'}]
    except Exception as e:
        print(f'Error parsing TXT {file_path}: {e}')
        return []


def parse_pptx(file_path: str) -> List[Dict]:
    segments = []
    if Presentation is None:
        return segments
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    slide_parts.append(shape.text)
            if slide_parts:
                segments.append({'text': f'Slide {i+1}\n' + '\n'.join(slide_parts), 'source_type': 'pptx_slide', 'page_num': i + 1})
    except Exception as e:
        print(f'Error parsing PPTX {file_path}: {e}')
    return segments


def parse_image_file(file_path: str) -> List[Dict]:
    try:
        image = Image.open(file_path)
        text = ocr_image(image)
        if text.strip():
            return [{'text': f'Image / engineering drawing OCR:\n{text}', 'source_type': 'image'}]
        return [{'text': 'Image uploaded. OCR could not extract readable text. This may still contain engineering drawings or pictures.', 'source_type': 'image'}]
    except Exception as e:
        print(f'Error parsing image {file_path}: {e}')
        return []


def parse_structured_text(file_path: str) -> List[Dict]:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return [{'text': f.read(), 'source_type': 'structured_text'}]
    except Exception as e:
        print(f'Error parsing structured text {file_path}: {e}')
        return []


def parse_document(file_path: str) -> List[Dict]:
    ext = file_path.lower().split('.')[-1]
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f'Unsupported file type: {ext}')
    if ext == 'pdf':
        return parse_pdf(file_path)
    if ext in ['doc', 'docx']:
        return parse_docx(file_path)
    if ext in ['xls', 'xlsx', 'csv']:
        return parse_excel_csv(file_path)
    if ext in ['txt', 'md', 'rtf', 'json', 'xml', 'html', 'htm']:
        return parse_txt(file_path)
    if ext in ['ppt', 'pptx']:
        return parse_pptx(file_path)
    if ext in ['png', 'jpg', 'jpeg', 'webp', 'bmp', 'tif', 'tiff']:
        return parse_image_file(file_path)
    raise ValueError(f'Unsupported file type: {ext}')


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    word_chunk_size = max(100, chunk_size // 5)
    word_overlap = min(word_chunk_size // 2, max(20, overlap // 5))
    while i < len(words):
        chunk_words = words[i:i + word_chunk_size]
        chunks.append(' '.join(chunk_words))
        if i + word_chunk_size >= len(words):
            break
        i += max(1, (word_chunk_size - word_overlap))
    return chunks


def process_file_into_chunks(file_path: str, plant: str, filename: str) -> List[Dict]:
    segments = parse_document(file_path)
    chunks = []
    for seg_idx, seg in enumerate(segments):
        raw_chunks = chunk_text(seg.get('text', ''), chunk_size=1000, overlap=200)
        for idx, c in enumerate(raw_chunks):
            if c.strip():
                chunks.append({
                    'chunk_id': f"{filename}_{seg.get('source_type', 'text')}_{seg_idx}_chunk_{idx}",
                    'plant': plant,
                    'filename': filename,
                    'filetype': filename.split('.')[-1].upper(),
                    'text': c,
                    'source_type': seg.get('source_type', 'text'),
                    'page_num': seg.get('page_num'),
                    'extra': seg.get('extra', {}),
                })
    return chunks
